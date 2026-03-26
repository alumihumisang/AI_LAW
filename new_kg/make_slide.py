"""生成複雜度比較投影片"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── 顏色定義 ──────────────────────────────────────────
DARK_BLUE   = RGBColor(0x1F, 0x39, 0x64)
MID_BLUE    = RGBColor(0x2E, 0x74, 0xB5)
LIGHT_BLUE  = RGBColor(0xD6, 0xE4, 0xF0)
ORANGE      = RGBColor(0xC5, 0x5A, 0x11)
LIGHT_ORANGE= RGBColor(0xFC, 0xE4, 0xD6)
GREEN       = RGBColor(0x37, 0x86, 0x44)
LIGHT_GREEN = RGBColor(0xD9, 0xEA, 0xD3)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
GRAY        = RGBColor(0xF2, 0xF2, 0xF2)
DARK_GRAY   = RGBColor(0x40, 0x40, 0x40)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

slide_layout = prs.slide_layouts[6]  # blank
slide = prs.slides.add_slide(slide_layout)

def add_rect(slide, l, t, w, h, fill_color=None, line_color=None, line_width=Pt(0)):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, text, l, t, w, h, font_size=14, bold=False, color=DARK_GRAY,
             align=PP_ALIGN.LEFT, italic=False):
    txBox = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return txBox

# ── 背景 ──────────────────────────────────────────────
add_rect(slide, 0, 0, 13.33, 7.5, fill_color=WHITE)

# ── 頂部藍色 banner ───────────────────────────────────
add_rect(slide, 0, 0, 13.33, 1.1, fill_color=DARK_BLUE)
add_text(slide, "Graph Construction Complexity Comparison",
         0.3, 0.1, 10, 0.55, font_size=28, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
add_text(slide, "O(N²) Pairwise  vs  XRAG Trie  ·  N = 6,057",
         0.3, 0.62, 10, 0.38, font_size=14, color=RGBColor(0xBD, 0xD7, 0xEE), align=PP_ALIGN.LEFT)

# ── 表格設定 ──────────────────────────────────────────
headers = ["", "O(N²) 逐對比較", "XRAG Trie"]
rows = [
    ["建構時間複雜度",   "O(N²)",              "O(N)"],
    ["每筆新案件成本",   "O(N)（重算 N 次相似度）", "O(1)（hash → 掛葉節點）"],
    ["Case 節點數",      "N",                  "N"],
    ["關係（邊）數",     "最多 N²/2 ≈ 1,830 萬", "N + ≤172 條 Trie 邊"],
    ["Neo4j Aura 上限",  "❌ 遠超 40 萬關係",   "✅ 完全在免費額度內"],
    ["實際建構時間",     "估計 > 51 小時",      "Phase2 數分鐘內完成"],
    ["新方法的代價",     "—",                  "Phase1 LLM 標注 ≈ 3–4 小時（一次性）"],
]

col_x = [0.35, 3.9, 8.55]
col_w = [3.4,  4.5,  4.5]
row_y_start = 1.3
row_h = 0.72

# header row
for ci, (hdr, cx, cw) in enumerate(zip(headers, col_x, col_w)):
    bg = DARK_BLUE if ci == 0 else (ORANGE if ci == 1 else MID_BLUE)
    add_rect(slide, cx, row_y_start, cw, 0.58, fill_color=bg,
             line_color=WHITE, line_width=Pt(1.5))
    add_text(slide, hdr, cx+0.08, row_y_start+0.06, cw-0.16, 0.46,
             font_size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# data rows
for ri, row in enumerate(rows):
    y = row_y_start + 0.58 + ri * row_h
    is_last = ri == len(rows) - 1
    for ci, (cell, cx, cw) in enumerate(zip(row, col_x, col_w)):
        # background
        if ci == 0:
            bg = RGBColor(0xE8, 0xEF, 0xF8)
        elif ri % 2 == 0:
            bg = GRAY
        else:
            bg = WHITE

        # highlight last row (trade-off)
        if is_last and ci == 2:
            bg = LIGHT_GREEN
        if "❌" in cell:
            bg = LIGHT_ORANGE
        if "✅" in cell:
            bg = LIGHT_GREEN

        add_rect(slide, cx, y, cw, row_h - 0.04,
                 fill_color=bg, line_color=RGBColor(0xCC, 0xCC, 0xCC), line_width=Pt(0.75))

        # text
        font_size = 12
        bold = (ci == 0)
        color = DARK_GRAY
        if "❌" in cell:
            color = ORANGE
        if "✅" in cell:
            color = GREEN
        if cell in ["O(N²)", "O(N)", "O(1)"]:
            bold = True
        add_text(slide, cell, cx+0.1, y+0.08, cw-0.2, row_h-0.15,
                 font_size=font_size, bold=bold, color=color, align=PP_ALIGN.CENTER)

# ── 底部 key insight box ──────────────────────────────
add_rect(slide, 0.35, 6.62, 12.65, 0.72,
         fill_color=RGBColor(0xE2, 0xEF, 0xDA),
         line_color=GREEN, line_width=Pt(1.5))
add_text(slide,
         "核心洞察：XRAG 將「比較問題（O(N²)）」轉為「分類問題（O(1)）」。"
         "代價是 Phase 1 的一次性 LLM 標注，換取之後每次 query 的閃電速度。",
         0.5, 6.65, 12.3, 0.65,
         font_size=12, bold=False, color=RGBColor(0x1E, 0x4D, 0x2B),
         align=PP_ALIGN.LEFT)

out = "/home/aru/AI_LAW/new_kg/complexity_comparison.pptx"
prs.save(out)
print(f"已儲存：{out}")
